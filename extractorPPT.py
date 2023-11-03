"""
Make sure to pip install python-pptx 
"""
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def extractPPT(pptFile, outputDir="Notes.md"):
    # Load the PowerPoint presentation
    
    presentation = Presentation(pptFile)
    text_file = outputDir
    # Create template
    file = open(text_file,"w",encoding="utf-8")
    file.write("<span style='color: #f2cf4a; font-family: Calibri;'></span>")
    file.close()


    # Create a subdirectory for images
    imgDirectory = "images/"+str(outputDir.split('.')[0])
    os.makedirs(imgDirectory, exist_ok=True)


    topLine = ""
    # Iterate through slides
    for slideNum, slide in enumerate(presentation.slides):
        with open(text_file, "a", encoding="utf-8") as file:
            #Add two newlines after slide
            
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if(slide.shapes.index(shape)==0 and len(slide.shapes)>1 and topLine!=str(paragraph.text)):
                            topLine=str(paragraph.text)
                            file.write("\n\n")
                            file.write(f"**{paragraph.text}:**\n")
                        else:
                            if(topLine==str(paragraph.text)): 
                                continue
                            #Add indentation based on paragraph level(for sub bullets).
                            indent = "".join("    " for _ in range(paragraph.level))
                            file.write(indent)
                            for run in paragraph.runs:
                                text = run.text
                                file.write(f"+ {text} ")
                            file.write("\n")
                
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:  # Check if shape is an image (shape_type ENUM 13)
                    imgFilePath = f"{imgDirectory}/img{slideNum}_{shape.shape_id}.png" #Generate image from slide
                    img = shape.image
                    imgByte = img.blob
                    with open(imgFilePath, "wb") as imgFile:
                        imgFile.write(imgByte)
                    file.write(f"![Image]({imgFilePath})\n")
    file.close()


if __name__ == "__main__":
    pptFile = "W8.pptx"  # Replace with your PowerPoint file
    outputDir = "W8.md"  # Replace with your desired output directory
    extractPPT(pptFile, outputDir)
